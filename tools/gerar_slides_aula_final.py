from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "AULA_FINAL_SPRING.pptx"
GUIDE = ROOT / "docs" / "ROTEIRO_SLIDES.md"

W = Inches(13.333)
H = Inches(7.5)

NAVY = "0B1220"
NAVY_2 = "111B2E"
BLUE = "2563EB"
CYAN = "22D3EE"
GREEN = "22C55E"
YELLOW = "FBBF24"
RED = "EF4444"
PURPLE = "8B5CF6"
WHITE = "F8FAFC"
MUTED = "94A3B8"
LINE = "334155"
CARD = "172033"


def rgb(value):
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def shape(slide, kind, x, y, w, h, fill=CARD, line=None, radius=False):
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line or fill)
    return shp


def text(slide, value, x, y, w, h, size=22, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def rich_text(slide, parts, x, y, w, h, size=22, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    for value, color, bold in parts:
        run = p.add_run()
        run.text = value
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def line(slide, x1, y1, x2, y2, color=LINE, width=2, arrow=False):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    if arrow:
        connector.line.end_arrowhead = True
    return connector


def pill(slide, value, x, y, w, color=BLUE, text_color=WHITE):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.42, color)
    text(slide, value, x, y, w, 0.42, 13, text_color, True, PP_ALIGN.CENTER)


def circle_label(slide, value, x, y, diameter, fill=BLUE, size=24):
    shape(slide, MSO_SHAPE.OVAL, x, y, diameter, diameter, fill)
    text(slide, value, x, y, diameter, diameter, size, WHITE, True, PP_ALIGN.CENTER)


def card(slide, title_value, body, x, y, w, h, accent=BLUE, icon=None):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, CARD, LINE)
    shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.07, h, accent, accent)
    if icon:
        circle_label(slide, icon, x + 0.25, y + 0.24, 0.56, accent, 17)
        tx = x + 0.95
    else:
        tx = x + 0.28
    text(slide, title_value, tx, y + 0.16, w - (tx - x) - 0.2, 0.46, 18, WHITE, True)
    text(slide, body, x + 0.28, y + 0.7, w - 0.56, h - 0.86, 15, MUTED, False,
         valign=MSO_ANCHOR.TOP)


def code(slide, value, x, y, w, h, accent=CYAN, size=16):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, "07101F", LINE)
    shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.07, h, accent, accent)
    text(slide, value, x + 0.24, y + 0.12, w - 0.42, h - 0.24, size, WHITE, False,
         font="Cascadia Code", valign=MSO_ANCHOR.TOP)


def base_slide(prs, title_value=None, number=None, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(NAVY)
    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.08, BLUE, BLUE)
    if section:
        pill(slide, section.upper(), 0.55, 0.32, max(1.25, len(section) * 0.105), PURPLE)
    if title_value:
        text(slide, title_value, 0.55, 0.7, 12.1, 0.72, 30, WHITE, True)
    if number is not None:
        text(slide, f"{number:02}", 12.35, 7.08, 0.45, 0.2, 10, MUTED, True, PP_ALIGN.RIGHT)
    return slide


def add_footer(slide, value="Aula final • Spring Boot + Thymeleaf"):
    text(slide, value, 0.55, 7.06, 7.0, 0.2, 9, MUTED)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 — capa
    s = base_slide(prs, number=1)
    shape(s, MSO_SHAPE.OVAL, 8.8, 0.8, 4.0, 4.0, "12284B", "1D4ED8")
    circle_label(s, "UI", 9.25, 1.32, 0.85, PURPLE, 18)
    circle_label(s, "SEC", 10.85, 1.05, 1.05, BLUE, 17)
    circle_label(s, "DB", 11.15, 2.72, 0.85, GREEN, 18)
    circle_label(s, "API", 9.35, 3.25, 1.0, YELLOW, 16)
    line(s, 10.05, 1.73, 10.85, 1.55, CYAN, 2)
    line(s, 11.35, 2.05, 11.55, 2.72, CYAN, 2)
    line(s, 10.35, 3.58, 11.15, 3.18, CYAN, 2)
    line(s, 9.72, 2.15, 9.75, 3.25, CYAN, 2)
    pill(s, "ÚLTIMO DIA", 0.72, 1.25, 1.55, RED)
    text(s, "Do formulário\nà aplicação segura", 0.72, 1.92, 7.6, 1.85, 42, WHITE, True)
    text(s, "Validação • DTO • Exceções • Security • REST", 0.75, 4.05, 7.7, 0.5, 20, CYAN, True)
    text(s, "Uma aula guiada por demonstrações", 0.75, 4.72, 6.4, 0.42, 17, MUTED)
    text(s, "Projeto Cardápio", 0.75, 6.55, 3.4, 0.35, 14, WHITE, True)
    add_footer(s, "Diego Alves • Aula final")

    # 2 — jornada
    s = base_slide(prs, "A jornada do projeto", 2, "retrospectiva")
    stages = [
        ("1", "MVC", "Controller + Thymeleaf", BLUE),
        ("2", "Banco", "JPA + MySQL", GREEN),
        ("3", "CRUD", "Incluir • Alterar • Excluir", YELLOW),
        ("4", "Login", "BCrypt + sessão", PURPLE),
        ("5", "Hoje", "Defesa em camadas", RED),
    ]
    x = 0.65
    for idx, (num, name, body, color) in enumerate(stages):
        circle_label(s, num, x, 2.0, 0.72, color, 20)
        text(s, name, x - 0.12, 2.9, 1.0, 0.34, 18, WHITE, True, PP_ALIGN.CENTER)
        text(s, body, x - 0.38, 3.35, 1.55, 0.75, 13, MUTED, False, PP_ALIGN.CENTER)
        if idx < len(stages) - 1:
            line(s, x + 0.75, 2.36, x + 2.34, 2.36, LINE, 3, True)
        x += 2.48
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.25, 5.0, 10.8, 0.88, "10213A", "1D4ED8")
    text(s, "Meta de hoje: transformar o CRUD funcional em uma aplicação que resiste a dados inválidos.",
         1.55, 5.12, 10.2, 0.55, 20, CYAN, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 3 — pergunta central
    s = base_slide(prs, "Podemos confiar no navegador?", 3, "problema")
    card(s, "O que o usuário vê", "Campo obrigatório\nPreço mínimo\nBotão Salvar\nModal amigável",
         0.7, 1.65, 4.65, 4.45, BLUE, "UI")
    card(s, "O que ele pode fazer", "Desativar JavaScript\nAlterar o HTML\nMontar outro POST\nModificar valores no console",
         7.98, 1.65, 4.65, 4.45, RED, "!")
    circle_label(s, "?", 5.91, 2.72, 1.5, YELLOW, 36)
    text(s, "O navegador está\nfora do nosso controle", 5.45, 4.46, 2.42, 0.9, 17, YELLOW, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 4 — demo
    s = base_slide(prs, "Demonstração: duas tentativas", 4, "live coding")
    text(s, "1", 0.75, 1.55, 0.45, 0.55, 30, GREEN, True)
    text(s, "Envio normal", 1.25, 1.55, 2.5, 0.55, 23, WHITE, True)
    code(s, "preço = -10\n\nsubmit → JavaScript\n       → modal\n       → envio cancelado", 0.75, 2.25, 5.25, 2.35, GREEN, 18)
    text(s, "Validação de experiência", 1.1, 4.86, 4.55, 0.4, 17, GREEN, True, PP_ALIGN.CENTER)

    text(s, "2", 7.05, 1.55, 0.45, 0.55, 30, RED, True)
    text(s, "Requisição adulterada", 7.55, 1.55, 3.6, 0.55, 23, WHITE, True)
    code(s, 'preco.value = "-10";\nform.submit();\n\nJavaScript ignorado\nbackend → erro.html', 7.05, 2.25, 5.25, 2.35, RED, 18)
    text(s, "Validação de segurança", 7.4, 4.86, 4.55, 0.4, 17, RED, True, PP_ALIGN.CENTER)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.45, 5.72, 6.45, 0.62, "2B1D0A", YELLOW)
    text(s, "Isto é manipulação da requisição — não SQL Injection.", 3.65, 5.8, 6.05, 0.4, 16, YELLOW, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 5 — camadas
    s = base_slide(prs, "Defesa em camadas", 5, "validação")
    layers = [
        ("JavaScript", "feedback imediato", PURPLE, 1.1, 5.5),
        ("DTO + @Valid", "contrato de entrada", BLUE, 1.75, 4.2),
        ("Service", "regra de negócio", GREEN, 2.4, 2.9),
        ("Banco", "integridade final", YELLOW, 3.05, 1.6),
    ]
    for label, sub, color, y, w in layers:
        x = (13.333 - w) / 2
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.9, color, color)
        text(s, label, x + 0.2, y + 0.05, w - 0.4, 0.34, 18, NAVY, True, PP_ALIGN.CENTER)
        text(s, sub, x + 0.2, y + 0.41, w - 0.4, 0.28, 12, NAVY, False, PP_ALIGN.CENTER)
    text(s, "Se uma camada for contornada, a próxima continua protegendo.",
         7.85, 2.2, 4.55, 1.2, 24, WHITE, True)
    text(s, "JS não substitui Java.\nJava não substitui o banco.",
         7.85, 3.65, 4.35, 0.9, 18, MUTED)
    add_footer(s)

    # 6 — DTO
    s = base_slide(prs, "DTO não é entidade", 6, "dto")
    card(s, "Tela / requisição", "nome\ndescrição\npreço", 0.7, 1.7, 3.0, 3.65, PURPLE, "1")
    card(s, "ItemCardapioDTO", "Somente os campos permitidos\n\n@NotBlank\n@DecimalMin", 5.15, 1.7, 3.0, 3.65, BLUE, "2")
    card(s, "Entidade JPA", "id\nnome\ndescrição\npreço\nincluidoPor", 9.6, 1.7, 3.0, 3.65, GREEN, "3")
    line(s, 3.75, 3.15, 5.05, 3.15, CYAN, 3, True)
    line(s, 8.2, 3.15, 9.5, 3.15, CYAN, 3, True)
    pill(s, "ENTRADA", 1.48, 5.72, 1.45, PURPLE)
    pill(s, "CONTRATO", 5.86, 5.72, 1.55, BLUE)
    pill(s, "BANCO", 10.45, 5.72, 1.3, GREEN)
    add_footer(s)

    # 7 — exceções
    s = base_slide(prs, "Uma saída única para os erros", 7, "exceções")
    error_cards = [
        ("Validação", "mensagem controlada", BLUE, 0.65),
        ("Item ausente", "mensagem do domínio", PURPLE, 3.75),
        ("Banco", "mensagem genérica", YELLOW, 6.85),
        ("Inesperado", "mensagem genérica", RED, 9.95),
    ]
    for title_value, body, color, x in error_cards:
        card(s, title_value, body, x, 1.55, 2.75, 1.55, color)
        line(s, x + 1.37, 3.15, 6.66, 4.05, color, 2, True)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 4.55, 4.0, 4.23, 1.05, "10213A", BLUE)
    text(s, "@ControllerAdvice", 4.75, 4.12, 3.83, 0.34, 22, CYAN, True, PP_ALIGN.CENTER)
    text(s, "escolhe a mensagem segura", 4.75, 4.5, 3.83, 0.28, 13, MUTED, False, PP_ALIGN.CENTER)
    line(s, 6.66, 5.08, 6.66, 5.65, CYAN, 3, True)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 5.15, 5.65, 3.03, 0.72, "341414", RED)
    text(s, "erro.html", 5.35, 5.75, 2.63, 0.35, 20, WHITE, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 8 — security before/after
    s = base_slide(prs, "Security limpa os controllers", 8, "spring security")
    pill(s, "ANTES", 0.8, 1.38, 1.2, RED)
    code(s, 'if (session.getAttribute(\n    "usuarioId") == null) {\n    return "redirect:/";\n}', 0.8, 2.0, 4.75, 2.05, RED, 17)
    text(s, "Repetido em cada método", 1.15, 4.35, 4.05, 0.42, 17, RED, True, PP_ALIGN.CENTER)
    circle_label(s, "→", 6.05, 2.55, 1.2, BLUE, 32)
    pill(s, "DEPOIS", 7.8, 1.38, 1.35, GREEN)
    code(s, '@GetMapping("/home")\npublic String mostrarHome(\n    Authentication auth, Model model) {\n    ...\n}', 7.8, 2.0, 4.75, 2.05, GREEN, 16)
    text(s, "O filtro barra antes do controller", 8.05, 4.35, 4.25, 0.42, 17, GREEN, True, PP_ALIGN.CENTER)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 2.65, 5.42, 8.05, 0.76, "10213A", BLUE)
    text(s, "Autenticação centralizada → menos repetição → menor chance de esquecer uma rota",
         2.9, 5.55, 7.55, 0.4, 16, CYAN, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 9 — login pipeline
    s = base_slide(prs, "O novo fluxo de login", 9, "spring security")
    items = [
        ("Form", "usuario + senha", PURPLE),
        ("Filtro", "intercepta /login", BLUE),
        ("Service", "busca por login", CYAN),
        ("BCrypt", "matches(senha, hash)", YELLOW),
        ("Context", "usuário autenticado", GREEN),
    ]
    x = 0.45
    for idx, (name, body, color) in enumerate(items):
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.05, 2.15, 1.55, CARD, color)
        shape(s, MSO_SHAPE.RECTANGLE, x, 2.05, 2.15, 0.1, color, color)
        text(s, name, x + 0.15, 2.34, 1.85, 0.35, 19, WHITE, True, PP_ALIGN.CENTER)
        text(s, body, x + 0.16, 2.82, 1.83, 0.45, 12, MUTED, False, PP_ALIGN.CENTER)
        if idx < len(items) - 1:
            line(s, x + 2.17, 2.82, x + 2.55, 2.82, color, 2, True)
        x += 2.62
    code(s, 'Authentication authentication\nauthentication.getName()', 4.05, 4.52, 5.25, 1.08, GREEN, 18)
    text(s, "O ID deixa de circular pelo formulário e pelos controllers.",
         3.6, 5.88, 6.15, 0.4, 17, GREEN, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 10 — csrf
    s = base_slide(prs, "CSRF: o selo secreto do formulário", 10, "segurança web")
    card(s, "Site legítimo", "Formulário Thymeleaf\n+ token CSRF correto", 0.8, 1.65, 3.4, 2.5, GREEN, "✓")
    circle_label(s, "TOKEN", 5.55, 2.04, 2.15, BLUE, 18)
    card(s, "Servidor", "Compara o token\nAceita ou rejeita o POST", 9.1, 1.65, 3.4, 2.5, BLUE, "S")
    line(s, 4.22, 2.9, 5.52, 2.9, GREEN, 3, True)
    line(s, 7.72, 2.9, 9.05, 2.9, GREEN, 3, True)
    card(s, "Outro site", "Tenta forjar o POST\nsem o token correto", 0.8, 4.65, 3.4, 1.45, RED, "×")
    line(s, 4.22, 5.35, 5.45, 4.2, RED, 3, True)
    shape(s, MSO_SHAPE.HEXAGON, 5.58, 4.07, 1.85, 1.7, "351515", RED)
    text(s, "403", 5.83, 4.47, 1.35, 0.42, 28, RED, True, PP_ALIGN.CENTER)
    text(s, "CSRF protege a origem do POST. Não valida nome, preço ou permissão de negócio.",
         8.25, 4.75, 4.15, 1.05, 18, YELLOW, True)
    add_footer(s)

    # 11 — mvc rest
    s = base_slide(prs, "MVC e REST: duas portas, um serviço", 11, "rest")
    card(s, "CardapioController", "@Controller\n\nreturn \"home\"\nHTML + Thymeleaf", 0.7, 1.55, 3.4, 3.65, PURPLE, "H")
    card(s, "CardapioService", "Regras de negócio\nValidação\nRepository", 4.96, 2.2, 3.4, 2.35, BLUE, "S")
    card(s, "CardapioRestController", "@RestController\n\nreturn DTO\nJSON", 9.23, 1.55, 3.4, 3.65, GREEN, "J")
    line(s, 4.15, 3.28, 4.91, 3.28, CYAN, 3, True)
    line(s, 8.41, 3.28, 9.18, 3.28, CYAN, 3, True)
    pill(s, "NAVEGADOR", 1.55, 5.62, 1.72, PURPLE)
    pill(s, "REUSO", 5.95, 5.12, 1.4, BLUE)
    pill(s, "API / MOBILE", 10.13, 5.62, 1.68, GREEN)
    text(s, "GET /api/cardapio", 4.82, 6.02, 3.7, 0.38, 18, CYAN, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 12 — arquitetura final
    s = base_slide(prs, "Arquitetura final", 12, "visão geral")
    nodes = [
        ("Browser", 0.55, 2.35, 1.55, PURPLE),
        ("Security", 2.65, 2.35, 1.65, RED),
        ("Controller", 4.85, 1.25, 1.75, BLUE),
        ("REST", 4.85, 3.55, 1.75, GREEN),
        ("DTO", 7.2, 2.35, 1.45, CYAN),
        ("Service", 9.2, 2.35, 1.55, YELLOW),
        ("Repository", 11.25, 1.25, 1.55, BLUE),
        ("MySQL", 11.25, 3.55, 1.55, GREEN),
    ]
    for name, x, y, w, color in nodes:
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.82, CARD, color)
        text(s, name, x + 0.08, y + 0.12, w - 0.16, 0.52, 16, WHITE, True, PP_ALIGN.CENTER)
    line(s, 2.12, 2.76, 2.6, 2.76, CYAN, 2, True)
    line(s, 4.32, 2.64, 4.8, 1.92, CYAN, 2, True)
    line(s, 4.32, 2.88, 4.8, 3.92, CYAN, 2, True)
    line(s, 6.65, 1.72, 7.15, 2.55, CYAN, 2, True)
    line(s, 6.65, 3.98, 7.15, 2.98, CYAN, 2, True)
    line(s, 8.68, 2.76, 9.15, 2.76, CYAN, 2, True)
    line(s, 10.78, 2.58, 11.2, 1.85, CYAN, 2, True)
    line(s, 12.02, 2.08, 12.02, 3.5, CYAN, 2, True)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, 5.35, 10.35, 0.82, "10213A", LINE)
    text(s, "Cada camada tem uma responsabilidade — e nenhuma confia cegamente na anterior.",
         1.75, 5.48, 9.85, 0.46, 20, CYAN, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 13 — roteiro aula
    s = base_slide(prs, "Condução da aula", 13, "roteiro")
    steps = [
        ("01", "CRUD funcionando", "5 min", BLUE),
        ("02", "Modal JavaScript", "10 min", PURPLE),
        ("03", "Contornar o JS", "15 min", RED),
        ("04", "DTO + validação", "25 min", CYAN),
        ("05", "Exceções", "20 min", YELLOW),
        ("06", "Spring Security", "30 min", GREEN),
        ("07", "REST + fechamento", "20 min", BLUE),
    ]
    y = 1.42
    for num, name, duration, color in steps:
        circle_label(s, num, 0.8, y, 0.55, color, 12)
        line(s, 1.38, y + 0.27, 2.0, y + 0.27, LINE, 2)
        text(s, name, 2.12, y - 0.02, 5.2, 0.42, 18, WHITE, True)
        pill(s, duration, 10.8, y + 0.02, 1.15, color, NAVY)
        y += 0.76
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 7.6, 2.05, 2.35, 2.75, "10213A", LINE)
    text(s, "RITMO", 7.9, 2.25, 1.75, 0.35, 16, CYAN, True, PP_ALIGN.CENTER)
    text(s, "conceito\n↓\ndemonstração\n↓\npergunta", 7.9, 2.78, 1.75, 1.55, 18, WHITE, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 14 — fechamento
    s = base_slide(prs, "O que eles precisam levar", 14, "fechamento")
    messages = [
        ("01", "O navegador não é confiável", RED),
        ("02", "DTO define o contrato de entrada", BLUE),
        ("03", "Service protege a regra de negócio", GREEN),
        ("04", "Security protege o fluxo", PURPLE),
        ("05", "MVC e REST podem compartilhar o núcleo", YELLOW),
    ]
    y = 1.45
    for num, message, color in messages:
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, y, 11.3, 0.72, CARD, LINE)
        circle_label(s, num, 1.2, y + 0.085, 0.55, color, 12)
        text(s, message, 2.05, y + 0.08, 9.65, 0.52, 19, WHITE, True)
        y += 0.92
    text(s, "Pergunta final", 1.0, 6.2, 2.0, 0.35, 15, CYAN, True)
    text(s, "Se eu remover uma camada, qual ataque ou erro passa a ser possível?",
         3.0, 6.05, 9.2, 0.6, 22, WHITE, True)
    add_footer(s)

    # 15 — fim
    s = base_slide(prs, number=15)
    pill(s, "ENCERRAMENTO", 0.75, 1.15, 1.65, GREEN)
    text(s, "Uma aplicação segura\nnão confia — verifica.", 0.75, 1.9, 8.2, 1.75, 42, WHITE, True)
    text(s, "JavaScript orienta. DTO limita. Service valida. Security protege.",
         0.78, 4.1, 10.5, 0.52, 21, CYAN, True)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.95, 1.45, 3.25, 3.25, "10213A", GREEN)
    circle_label(s, "✓", 9.9, 2.05, 1.35, GREEN, 36)
    text(s, "PROJETO\nCOMPLETO", 9.35, 3.55, 2.45, 0.72, 17, WHITE, True, PP_ALIGN.CENTER)
    text(s, "Perguntas?", 0.78, 5.65, 3.2, 0.55, 28, YELLOW, True)
    add_footer(s, "Obrigado!")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)

    guide = """# Roteiro do apresentador — Aula final

## Como usar

- Apresentação: `AULA_FINAL_SPRING.pptx`
- Duração sugerida: 2h a 2h30.
- Regra de ritmo: no máximo dois slides seguidos sem abrir o projeto.

## Slide 1 — Do formulário à aplicação segura

Abra dizendo que a aula não adiciona apenas funcionalidades: ela transforma o CRUD em uma aplicação capaz de rejeitar entradas maliciosas ou incorretas.

## Slide 2 — A jornada

Recapitule rapidamente MVC, banco, CRUD e login. Pergunte aos alunos onde eles acham que ainda existe fragilidade.

## Slide 3 — Podemos confiar no navegador?

Mostre que `required`, `min` e JavaScript estão na máquina do usuário. Ele controla o navegador; nós controlamos somente o servidor.

## Slide 4 — Demonstração

1. Tente salvar preço `-10` normalmente: mostre o modal.
2. Abra o console do navegador.
3. Execute:

```javascript
document.querySelector("#preco").value = "-10";
document.querySelector("#form-cardapio").submit();
```

4. Mostre `erro.html` vindo do backend.
5. Reforce: isso é manipulação da requisição, não SQL Injection.

## Slide 5 — Defesa em camadas

Explique de fora para dentro. JavaScript melhora UX; DTO valida contrato; service garante negócio; banco mantém integridade.

## Slide 6 — DTO

Abra `ItemCardapioDTO`. Compare com `ItemCardapio`: o DTO não contém `incluidoPor`. Isso impede que o cliente escolha quem criou o item.

## Slide 7 — Exceções

Abra `TratadorDeExcecoes`. Mostre por que validação pode exibir mensagem controlada, mas banco e erro inesperado recebem texto genérico.

## Slide 8 — Security limpa controllers

Compare mentalmente o antigo `if (session...)` repetido com os controllers atuais. A proteção ocorre antes do controller.

## Slide 9 — Login

Abra `SecurityConfig` e `LoginService`. Siga a sequência do slide. Mostre `authentication.getName()` no controller.

## Slide 10 — CSRF

Inspecione o HTML renderizado e procure o campo `_csrf`. Explique que `th:action` permite a inclusão automática do token nos formulários POST.

## Slide 11 — MVC x REST

Depois de autenticar, abra `http://localhost:8080/api/cardapio`. Compare JSON com `home.html`. Mostre que os dois controllers usam o mesmo service.

## Slide 12 — Arquitetura final

Faça os alunos narrarem o caminho de uma inclusão: browser, Security, controller, DTO, service, repository e MySQL.

## Slide 13 — Condução

Use como checkpoint de tempo. Se estiver atrasado, reduza a parte REST, mas preserve a demonstração de validação e o Security.

## Slide 14 — O que levar

Peça um exemplo concreto para cada frase. Exemplo: “Como alguém contorna o JavaScript?”

## Slide 15 — Encerramento

Feche com a frase: “Uma aplicação segura não confia — verifica”. Abra para perguntas.
"""
    GUIDE.write_text(guide, encoding="utf-8")

    print(f"PPTX: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Guia: {GUIDE}")


if __name__ == "__main__":
    build()
